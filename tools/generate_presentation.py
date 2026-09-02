"""
Генератор корпоративной презентации PowerPoint (.pptx) для IntraLink AI.
Создает 15 слайдов в формате 16:9 с темной темой, карточками, метриками,
скриншотами реального интерфейса и полными заметками докладчика.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- Цветовая палитра IntraLink Executive Enterprise ---
COLOR_BG = RGBColor(12, 13, 14)          # #0C0D0E (Основной фон)
COLOR_CARD_BG = RGBColor(20, 23, 26)     # #14171A (Поверхность карточек)
COLOR_CARD_BORDER = RGBColor(40, 46, 54) # #282E36 (Границы)
COLOR_TEXT_PRIMARY = RGBColor(248, 250, 252) # #F8FAFC (Белый текст)
COLOR_TEXT_SECONDARY = RGBColor(148, 163, 184) # #94A3B8 (Серый подзаголовок)
COLOR_TEXT_MUTED = RGBColor(100, 116, 139) # #64748B (Мелкий текст)

# Акцентные цвета
COLOR_ACCENT_BLUE = RGBColor(59, 130, 246)   # #3B82F6 (AI / Интеллект)
COLOR_ACCENT_EMERALD = RGBColor(16, 185, 129) # #10B981 (ROI / Успех / SLA)
COLOR_ACCENT_AMBER = RGBColor(245, 158, 11)  # #F59E0B (Внимание / Дубликаты)
COLOR_ACCENT_ROSE = RGBColor(244, 63, 94)    # #F43F5E (Проблемы / Рутина)
COLOR_ACCENT_PURPLE = RGBColor(139, 92, 246) # #8B5CF6 (RAG / Векторы)

FONT_MAIN = "Inter"
FONT_MONO = "JetBrains Mono"

SCREENSHOTS_DIR = os.path.join("docs", "presentations", "screenshots")

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    def add_bg(slide):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = COLOR_BG
        bg.line.fill.background()
        return bg

    def add_header(slide, category, title, subtitle=""):
        # Категория / Бейдж
        tb_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.32))
        tf_cat = tb_cat.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category.upper()
        p_cat.font.name = FONT_MONO
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_ACCENT_BLUE

        # Основной заголовок
        tb_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.68), Inches(11.7), Inches(0.55))
        tf_title = tb_title.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_top = tf_title.margin_right = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title
        p_title.font.name = FONT_MAIN
        p_title.font.size = Pt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_TEXT_PRIMARY

        # Подзаголовок
        if subtitle:
            tb_sub = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(11.7), Inches(0.35))
            tf_sub = tb_sub.text_frame
            tf_sub.word_wrap = True
            tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.name = FONT_MAIN
            p_sub.font.size = Pt(11.5)
            p_sub.font.color.rgb = COLOR_TEXT_SECONDARY

    def add_card(slide, left, top, width, height, title="", border_color=COLOR_CARD_BORDER, bg_color=COLOR_CARD_BG):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = border_color
        card.line.width = Pt(1.2)
        
        if title:
            tb = slide.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.18), Inches(width - 0.5), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.name = FONT_MAIN
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = COLOR_TEXT_PRIMARY
        return card

    def add_screenshot_frame(slide, left, top, width, height, image_filename, caption=""):
        # Карточка подложка с границей
        frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
        frame.fill.solid()
        frame.fill.fore_color.rgb = COLOR_CARD_BG
        frame.line.color.rgb = COLOR_CARD_BORDER
        frame.line.width = Pt(1)

        img_path = os.path.join(SCREENSHOTS_DIR, image_filename)
        if os.path.exists(img_path):
            img_margin = 0.08
            img_top = top + img_margin
            img_left = left + img_margin
            img_w = width - (img_margin * 2)
            img_h = height - (0.4 if caption else (img_margin * 2))
            
            slide.shapes.add_picture(img_path, Inches(img_left), Inches(img_top), width=Inches(img_w), height=Inches(img_h))
            
            if caption:
                tb_c = slide.shapes.add_textbox(Inches(left + 0.15), Inches(top + height - 0.32), Inches(width - 0.3), Inches(0.25))
                tf_c = tb_c.text_frame
                tf_c.word_wrap = True
                tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
                p_c = tf_c.paragraphs[0]
                p_c.text = caption
                p_c.font.name = FONT_MONO
                p_c.font.size = Pt(9.5)
                p_c.font.color.rgb = COLOR_TEXT_MUTED
        return frame

    # =========================================================================
    # СЛАЙД 1: Титульный (Title Slide)
    # =========================================================================
    s1 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s1)
    
    badge = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.6), Inches(3.6), Inches(0.42))
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(19, 36, 64)
    badge.line.color.rgb = COLOR_ACCENT_BLUE
    badge.line.width = Pt(1)
    tf_b = badge.text_frame
    p_b = tf_b.paragraphs[0]
    p_b.text = "NEXT-GEN IT HELPDESK & AI HUB"
    p_b.alignment = PP_ALIGN.CENTER
    p_b.font.name = FONT_MONO
    p_b.font.size = Pt(10.5)
    p_b.font.bold = True
    p_b.font.color.rgb = COLOR_ACCENT_BLUE

    tb_main = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.6))
    tf_main = tb_main.text_frame
    tf_main.word_wrap = True
    p_m1 = tf_main.paragraphs[0]
    p_m1.text = "IntraLink AI"
    p_m1.font.name = FONT_MAIN
    p_m1.font.size = Pt(44)
    p_m1.font.bold = True
    p_m1.font.color.rgb = COLOR_TEXT_PRIMARY

    p_m2 = tf_main.add_paragraph()
    p_m2.text = "Интеллектуальная трансформация Helpdesk 1-й линии"
    p_m2.font.name = FONT_MAIN
    p_m2.font.size = Pt(24)
    p_m2.font.bold = True
    p_m2.font.color.rgb = COLOR_ACCENT_BLUE

    tb_desc = s1.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.0), Inches(1.0))
    tf_desc = tb_desc.text_frame
    tf_desc.word_wrap = True
    p_d = tf_desc.paragraphs[0]
    p_d.text = "Автоматизация триажа, семантическая база знаний RAG (pgvector), локальная LLM-суммаризация, компьютерное зрение и роботизация инфраструктурных задач в закрытом контуре компании."
    p_d.font.name = FONT_MAIN
    p_d.font.size = Pt(13.5)
    p_d.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s1, 0.8, 5.5, 3.6, 1.3, border_color=COLOR_ACCENT_EMERALD)
    tb_s1 = s1.shapes.add_textbox(Inches(1.0), Inches(5.65), Inches(3.2), Inches(0.9))
    tf_s1 = tb_s1.text_frame
    p_s1_v = tf_s1.paragraphs[0]
    p_s1_v.text = "-94% MTTR"
    p_s1_v.font.name = FONT_MAIN
    p_s1_v.font.size = Pt(20)
    p_s1_v.font.bold = True
    p_s1_v.font.color.rgb = COLOR_ACCENT_EMERALD
    p_s1_l = tf_s1.add_paragraph()
    p_s1_l.text = "Сокращение времени решения инцидентов"
    p_s1_l.font.size = Pt(11)
    p_s1_l.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s1, 4.8, 5.5, 3.6, 1.3, border_color=COLOR_ACCENT_BLUE)
    tb_s2 = s1.shapes.add_textbox(Inches(5.0), Inches(5.65), Inches(3.2), Inches(0.9))
    tf_s2 = tb_s2.text_frame
    p_s2_v = tf_s2.paragraphs[0]
    p_s2_v.text = "100% On-Premise"
    p_s2_v.font.name = FONT_MAIN
    p_s2_v.font.size = Pt(20)
    p_s2_v.font.bold = True
    p_s2_v.font.color.rgb = COLOR_ACCENT_BLUE
    p_s2_l = tf_s2.add_paragraph()
    p_s2_l.text = "Полная безопасность и закрытый контур"
    p_s2_l.font.size = Pt(11)
    p_s2_l.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s1, 8.8, 5.5, 3.7, 1.3, border_color=COLOR_ACCENT_PURPLE)
    tb_s3 = s1.shapes.add_textbox(Inches(9.0), Inches(5.65), Inches(3.3), Inches(0.9))
    tf_s3 = tb_s3.text_frame
    p_s3_v = tf_s3.paragraphs[0]
    p_s3_v.text = "160+ ч/мес"
    p_s3_v.font.name = FONT_MAIN
    p_s3_v.font.size = Pt(20)
    p_s3_v.font.bold = True
    p_s3_v.font.color.rgb = COLOR_ACCENT_PURPLE
    p_s3_l = tf_s3.add_paragraph()
    p_s3_l.text = "Экономия рабочего времени инженеров"
    p_s3_l.font.size = Pt(11)
    p_s3_l.font.color.rgb = COLOR_TEXT_SECONDARY

    s1.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Добрый день, коллеги и руководство! Сегодня я представляю результаты внедрения интеллектуальной "
        "системы IntraLink AI. Наша цель — кардинальное ускорение 1-й линии поддержки, снижение рутины на инженеров "
        "и повышение соблюдения SLA с помощью современных AI-технологий, работающих строго внутри нашего закрытого корпоративного периметра»."
    )

    # =========================================================================
    # СЛАЙД 2: Проблематика (Challenges)
    # =========================================================================
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s2)
    add_header(s2, "Текущие вызовы", "Узкие места традиционной 1-й линии поддержки", "С чем сталкивается ИТ-отдел при классическом ручном обслуживании заявок")

    col_w = 3.65
    top_pos = 1.8
    add_card(s2, 0.8, top_pos, col_w, 5.1, "1. Высокая доля рутины", border_color=COLOR_ACCENT_ROSE)
    tb_c1 = s2.shapes.add_textbox(Inches(1.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_c1 = tb_c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.paragraphs[0].text = "• До 40-50% заявок — однотипные задачи (Wi-Fi, сброс паролей, подключение принтеров, консультации)."
    tf_c1.paragraphs[0].font.size = Pt(11.5)
    tf_c1.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_c1.add_paragraph()
    p.text = "• Инженер тратит до 15 минут только на открытие, смену статуса, вбивание шаблона и закрытие каждого тикета."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_c1.add_paragraph()
    p.text = "• В пиковые часы накапливается очередь («завал»), срываются регламентные нормативы SLA."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s2, 4.8, top_pos, col_w, 5.1, "2. Потеря накопленного опыта", border_color=COLOR_ACCENT_AMBER)
    tb_c2 = s2.shapes.add_textbox(Inches(5.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_c2 = tb_c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.paragraphs[0].text = "• Исторические решения заперты в архиве (10 000+ заявок без удобного семантического поиска)."
    tf_c2.paragraphs[0].font.size = Pt(11.5)
    tf_c2.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_c2.add_paragraph()
    p.text = "• Новый инженер заново исследует сбой, который уже решался коллегой на прошлой неделе."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_c2.add_paragraph()
    p.text = "• Рост MTTR (времени решения) сложных инцидентов из-за долгого ручного поиска инструкций."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s2, 8.8, top_pos, col_w, 5.1, "3. Дубли и мусорный поток", border_color=COLOR_ACCENT_BLUE)
    tb_c3 = s2.shapes.add_textbox(Inches(9.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_c3 = tb_c3.text_frame
    tf_c3.word_wrap = True
    tf_c3.paragraphs[0].text = "• До 25% входящего потока — повторные заявки-дубликаты от нетерпеливых заявителей."
    tf_c3.paragraphs[0].font.size = Pt(11.5)
    tf_c3.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_c3.add_paragraph()
    p.text = "• Ошибочные направления: заявки в ИТ, предназначенные для АХО, безопасности или кадров."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_c3.add_paragraph()
    p.text = "• Ручная отмена и разъяснения отнимают часы у квалифицированных специалистов."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    s2.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«В традиционном Helpdesk до половины времени инженера уходит на монотонные операции: "
        "войти в тикет, сменить исполнителя, скопировать логин, пойти в AD, закрыть. При этом 10 000 архивных кейсов "
        "лежат мертвым грузом, а очередь засорена дубликатами. IntraLink AI устраняет эти проблемы»."
    )

    # =========================================================================
    # СЛАЙД 3: Архитектура (Architecture)
    # =========================================================================
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s3)
    add_header(s3, "Архитектура решения", "Комплексная AI-экосистема IntraLink", "Единый интеллектуальный мозг, детерминированные правила и прямое исполнение")

    add_card(s3, 0.8, 1.8, 11.7, 1.5, "1. Слой восприятия и интерфейсов (Интерактивный ввод)", border_color=COLOR_ACCENT_BLUE)
    tb_l1 = s3.shapes.add_textbox(Inches(1.05), Inches(2.35), Inches(11.2), Inches(0.8))
    tf_l1 = tb_l1.text_frame
    tf_l1.word_wrap = True
    p = tf_l1.paragraphs[0]
    p.text = "• Telegram-бот: Моментальные Push-оповещения за 0.3 сек через Redis Streams с Consumer Groups.\n• Web SPA Cockpit (/admin): Интерактивный диспетчерский пульт с инспектором заявок, фильтрами и канбаном.\n• CLI Execution Agent: Пакетный разбор и диагностика через слэш-команды."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s3, 0.8, 3.5, 11.7, 1.8, "2. Интеллектуальное ядро Core API & AI Hub", border_color=COLOR_ACCENT_PURPLE)
    tb_l2 = s3.shapes.add_textbox(Inches(1.05), Inches(4.05), Inches(11.2), Inches(1.1))
    tf_l2 = tb_l2.text_frame
    tf_l2.word_wrap = True
    p = tf_l2.paragraphs[0]
    p.text = "• Rule Engine: Детерминированный триаж, классификация и расчет жестких SLA.\n• RAG Engine: Семантический векторный поиск решений в PostgreSQL 16 (pgvector + FastEmbed).\n• Local LLM Engine: Локальная нейросеть Ollama (Qwen2.5:1.5B) для моментальной суммаризации переписки.\n• Vision OCR: Мультимодальный анализ скриншотов сбоев, ошибок 1С и BSOD."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s3, 0.8, 5.5, 11.7, 1.4, "3. Уровень исполнения в инфраструктуре (Zero-Click Execution)", border_color=COLOR_ACCENT_EMERALD)
    tb_l3 = s3.shapes.add_textbox(Inches(1.05), Inches(6.0), Inches(11.2), Inches(0.75))
    tf_l3 = tb_l3.text_frame
    tf_l3.word_wrap = True
    p = tf_l3.paragraphs[0]
    p.text = "• Active Directory PowerShell (WLAN-WORKNET / Учетные записи)  |  WinRM / SMB (Установка принтеров)\n• Сетевая экспресс-диагностика (Ping, DNS, SMB:445, WinRM:5985) с защитой Verified Execution."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    s3.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Архитектура IntraLink объединяет три слоя: сверху — мгновенные интерфейсы (Telegram и Web Cockpit), "
        "в центре — Core API с базой pgvector и локальной LLM, внизу — прямые защищенные модули исполнения "
        "в Active Directory и Windows-домене. Тикет переходит в статус 'Выполнена' только после физического подтверждения от контроллера»."
    )

    # =========================================================================
    # СЛАЙД 4: Демонстрация: Диспетчерский пульт (Live Showcase 1)
    # =========================================================================
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s4)
    add_header(s4, "Демонстрация системы", "Интеллектуальный реестр и диспетчерский пульт", "Оперативный контроль потока заявок в реальном времени со смарт-фильтрацией")

    add_card(s4, 0.8, 1.8, 4.8, 5.1, "Возможности пульта", border_color=COLOR_ACCENT_BLUE)
    tb_p4 = s4.shapes.add_textbox(Inches(1.05), Inches(2.4), Inches(4.3), Inches(4.3))
    tf_p4 = tb_p4.text_frame
    tf_p4.word_wrap = True
    p = tf_p4.paragraphs[0]
    p.text = "• Смарт-фильтры по 11 направлениям:\nВсе, Wi-Fi, 1C, Доступы, Принтеры, Ремонт (Каб. 112), Сеть, Почта, Оборудование."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p4.add_paragraph()
    p.text = "• Цветовая шкала SLA:\nМгновенная подсветка тикетов с приближающимся дедлайном (Critical / Urgent / Normal)."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p4.add_paragraph()
    p.text = "• Автопривязка ПК и контактов:\nИзвлечение имени рабочей станции (`NB-0081`, `PC-114`) и контактов заявителя."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_screenshot_frame(s4, 5.8, 1.8, 6.7, 5.1, "01_queue_dashboard.png", "Реальный вид очереди заявок Helpdesk с активными фильтрами и статусами")

    s4.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«На экране — рабочий диспетчерский пульт IntraLink. Инженер видит очередь заявок с автоматической "
        "категоризацией по направлениям, таймерами SLA и извлеченными именами компьютеров. "
        "Всё обновляется в реальном времени без перезагрузки страниц»."
    )

    # =========================================================================
    # СЛАЙД 5: AI Фича №1 — RAG База знаний
    # =========================================================================
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s5)
    add_header(s5, "AI Feature 01", "Семантический RAG-поиск в базе решений", "Мгновенное извлечение готовых решений из 10 000+ архивных кейсов")

    add_card(s5, 0.8, 1.8, 5.7, 5.1, "Как работает технология RAG", border_color=COLOR_ACCENT_PURPLE)
    tb_r1 = s5.shapes.add_textbox(Inches(1.05), Inches(2.45), Inches(5.2), Inches(4.2))
    tf_r1 = tb_r1.text_frame
    tf_r1.word_wrap = True
    p = tf_r1.paragraphs[0]
    p.text = "1. Автосинхронизация (`/sync`):\nКаждая решенная инженером заявка векторизуется моделью FastEmbed и сохраняется в pgvector."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_r1.add_paragraph()
    p.text = "2. Семантическое понимание сути:\nПоиск находит релевантные решения не по точным ключевым словам, а по смыслу проблемы (синонимы, опечатки, разное описание сбоя пользователем)."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_r1.add_paragraph()
    p.text = "3. Время отклика:\nПоиск по 10 000+ историческим кейсам занимает менее 50 миллисекунд."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s5, 6.8, 1.8, 5.7, 5.1, "Бизнес-эффект и выгода", border_color=COLOR_ACCENT_EMERALD)
    tb_r2 = s5.shapes.add_textbox(Inches(7.05), Inches(2.45), Inches(5.2), Inches(4.2))
    tf_r2 = tb_r2.text_frame
    tf_r2.word_wrap = True
    p = tf_r2.paragraphs[0]
    p.text = "• Сокращение времени диагностики сложных инцидентов с 30 минут до 1 минуты."
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_EMERALD
    p = tf_r2.add_paragraph()
    p.text = "• Быстрый онбординг новичков: начинающий инженер сразу видит точные шаги решения от ведущих экспертов."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_r2.add_paragraph()
    p.text = "• Стандартизация качества ответов: заявители получают проверенные инструкции."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_r2.add_paragraph()
    p.text = "• Zero Cloud Leakage: все векторные базы и модели работают строго на корпоративном сервере."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    s5.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Первая ключевая AI-фича — RAG на базе pgvector. Раньше опыт уходил вместе с сотрудниками. "
        "Теперь система векторизует закрытые заявки. При поступлении новой проблемы инженер за 50 миллисекунд "
        "получает 3 лучших исторических кейса решения с готовыми инструкциями»."
    )

    # =========================================================================
    # СЛАЙД 6: AI Фича №2 — Инспектор заявки и Автоответы (Live Showcase 2)
    # =========================================================================
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s6)
    add_header(s6, "AI Feature 02", "Инспектор тикета и генерация черновиков", "Контекстный анализ заявки, обнаружение дубликатов и персонализированные ответы")

    add_card(s6, 0.8, 1.8, 4.8, 5.1, "Интеллект инспектора", border_color=COLOR_ACCENT_PURPLE)
    tb_p6 = s6.shapes.add_textbox(Inches(1.05), Inches(2.4), Inches(4.3), Inches(4.3))
    tf_p6 = tb_p6.text_frame
    tf_p6.word_wrap = True
    p = tf_p6.paragraphs[0]
    p.text = "• Авто-генерация ответов:\nСистема анализирует заявку и формирует персонализированный вежливый черновик ответа с подстановкой имени заявителя и регламентных шагов."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p6.add_paragraph()
    p.text = "• Авто-определение дубликата:\nИнспектор подсвечивает: «Обнаружен дубликат заявки #139929» с кнопкой объединения."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p6.add_paragraph()
    p.text = "• Режимы отправки:\nОтправка комментария, перевод в ожидание или автоматическое закрытие с фиксацией трудозатрат."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_screenshot_frame(s6, 5.8, 1.8, 6.7, 5.1, "04_auto_reply_ai.png", "Инспектор заявки: AI черновик ответа и предупреждение о дубликате")

    s6.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«На скриншоте видно работу AI-инспектора. Система сама поняла суть обращения, "
        "сгенерировала корректный ответ заявителю и сразу предупредила инженера о наличии дубликата "
        "тикета #139929. Инженеру остается только нажать кнопку подтверждения»."
    )

    # =========================================================================
    # СЛАЙД 7: AI Фича №3 — Умные дубликаты и авто-редирект
    # =========================================================================
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s7)
    add_header(s7, "AI Feature 03", "Фильтрация шума: Умные дубликаты и авто-редирект", "Автоматическая очистка очереди от повторов и ошибочно направленных тикетов")

    add_card(s7, 0.8, 1.8, 5.7, 5.1, "Автоматическое слияние дубликатов", border_color=COLOR_ACCENT_AMBER)
    tb_d1 = s7.shapes.add_textbox(Inches(1.05), Inches(2.45), Inches(5.2), Inches(4.2))
    tf_d1 = tb_d1.text_frame
    tf_d1.word_wrap = True
    p = tf_d1.paragraphs[0]
    p.text = "• Проблема:\nПользователи создают 2-3 одинаковые заявки подряд (через портал, почту и телефон)."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_d1.add_paragraph()
    p.text = "• AI-Решение (`/duplicates`):\nАлгоритм находит семантические и временные совпадения, определяет Master Ticket (основную заявку), переводит копии в статус «Отменена» и оставляет пользователю вежливый комментарий со ссылкой на основной тикет."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_d1.add_paragraph()
    p.text = "• Эффект: Полное исключение двойной работы разными инженерами."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_AMBER

    add_card(s7, 6.8, 1.8, 5.7, 5.1, "Интеллектуальный редирект сервисов", border_color=COLOR_ACCENT_BLUE)
    tb_d2 = s7.shapes.add_textbox(Inches(7.05), Inches(2.45), Inches(5.2), Inches(4.2))
    tf_d2 = tb_d2.text_frame
    tf_d2.word_wrap = True
    p = tf_d2.paragraphs[0]
    p.text = "• Проблема:\nЗаявки на ремонт мебели, пропуска или замену ламп по ошибке падают в очередь ИТ-отдела."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_d2.add_paragraph()
    p.text = "• AI-Решение (`/redirect`):\nСистема распознает нецелевой сервис, отменяет ошибочный тикет в ИТ-очереди и отправляет заявителю точную ссылку на нужную услугу в каталоге IntraService (АХО, СБ, Кадры)."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_d2.add_paragraph()
    p.text = "• Эффект: Очередь ИТ чиста от непрофильных задач."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_BLUE

    s7.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Около 25-30% заявок в корпоративной поддержке — это повторные дубли либо заявки не по адресу. "
        "IntraLink автоматически связывает дубли с основным тикетом и отменяет их с понятным комментарием. "
        "А заявки в АХО или СБ мгновенно перенаправляются по нужным ссылкам без отвлечения инженеров»."
    )

    # =========================================================================
    # СЛАЙД 8: AI Фича №4 — Сетевая экспресс-диагностика (Live Showcase 3)
    # =========================================================================
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s8)
    add_header(s8, "AI Feature 04", "Экспресс-диагностика сети и рабочих станций", "Мгновенная проверка доступности хоста заявителя по сетевым протоколам")

    add_card(s8, 0.8, 1.8, 4.8, 5.1, "Сетевой сканер (/diag)", border_color=COLOR_ACCENT_BLUE)
    tb_p8 = s8.shapes.add_textbox(Inches(1.05), Inches(2.4), Inches(4.3), Inches(4.3))
    tf_p8 = tb_p8.text_frame
    tf_p8.word_wrap = True
    p = tf_p8.paragraphs[0]
    p.text = "• Параллельная 4-точечная проверка:\n1. ICMP Ping (базовая связность)\n2. DNS Resolve (соответствие IP и имени хоста)\n3. SMB Port 445 (файловый доступ)\n4. WinRM Port 5985 (удаленное управление)."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p8.add_paragraph()
    p.text = "• Скорость выполнения:\nВсе проверки выполняются параллельно менее чем за 1 секунду."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p8.add_paragraph()
    p.text = "• Практическая польза:\nИнженер еще до звонка заявителю знает, включен ли компьютер и готов ли он к удаленному подключению."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_screenshot_frame(s8, 5.8, 1.8, 6.7, 5.1, "03_network_diagnostics.png", "Результат сетевой экспресс-диагностики рабочего места заявителя")

    s8.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«На скриншоте показан модуль сетевой экспресс-диагностики. "
        "При нажатии кнопки 'Диагностика сети' система за секунду проверяет Ping, DNS, порт SMB и WinRM. "
        "Инженер сразу видит, доступен ли ПК пользователя, не тратя время на слепые попытки подключения»."
    )

    # =========================================================================
    # СЛАЙД 9: AI Фича №5 — Компьютерное зрение и локальная LLM
    # =========================================================================
    s9 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s9)
    add_header(s9, "AI Feature 05", "Компьютерное зрение и локальная нейросеть", "Мультимодальный анализ скриншотов и моментальная суммаризация диалогов")

    add_card(s9, 0.8, 1.8, 5.7, 5.1, "Vision OCR: Анализ скриншотов (/screen)", border_color=COLOR_ACCENT_BLUE)
    tb_v1 = s9.shapes.add_textbox(Inches(1.05), Inches(2.45), Inches(5.2), Inches(4.2))
    tf_v1 = tb_v1.text_frame
    tf_v1.word_wrap = True
    p = tf_v1.paragraphs[0]
    p.text = "• Автоскачивание вложений из IntraService."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_v1.add_paragraph()
    p.text = "• Распознавание системных окон: ошибки 1С, синие экраны (BSOD), сетевые ошибки Windows, сбои драйверов."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_v1.add_paragraph()
    p.text = "• Формирование гипотезы первопричины (Root Cause) и чек-листа решения до первого контакта с пользователем."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s9, 6.8, 1.8, 5.7, 5.1, "Локальная LLM: Суммаризация (Ollama)", border_color=COLOR_ACCENT_PURPLE)
    tb_v2 = s9.shapes.add_textbox(Inches(7.05), Inches(2.45), Inches(5.2), Inches(4.2))
    tf_v2 = tb_v2.text_frame
    tf_v2.word_wrap = True
    p = tf_v2.paragraphs[0]
    p.text = "• Модель: Qwen2.5:1.5B в изолированном Docker-контейнере."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_v2.add_paragraph()
    p.text = "• Сжатие переписки: Анализирует длинную ветку из 10+ комментариев за 1.3 секунды."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_v2.add_paragraph()
    p.text = "• Структурированная выжимка:\n1. Суть проблемы\n2. Что уже предпринималось\n3. Текущий блокер и следующий шаг."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_v2.add_paragraph()
    p.text = "• 100% приватность: Никакие данные не отправляются во внешние облака."
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_PURPLE

    s9.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Четвертая и пятая фичи — мультимодальный стек. Модуль /screen автоматически считывает текст ошибки "
        "со скриншотов пользователей. А встроенная модель Qwen2.5 сжимает длинные переписки в три четких пункта "
        "всего за 1.3 секунды, экономя время при пересменке дежурных»."
    )

    # =========================================================================
    # СЛАЙД 10: AI Фича №6 — Роботизация исполнения в инфраструктуре
    # =========================================================================
    s10 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s10)
    add_header(s10, "AI Feature 06", "Роботизация инфраструктуры (Zero-Click Execution)", "Автоматическое применение изменений в Active Directory и Windows-домене")

    col_w = 3.65
    top_pos = 1.8
    add_card(s10, 0.8, top_pos, col_w, 5.1, "1. Active Directory (Wi-Fi / УЗ)", border_color=COLOR_ACCENT_EMERALD)
    tb_e1 = s10.shapes.add_textbox(Inches(1.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_e1 = tb_e1.text_frame
    tf_e1.word_wrap = True
    tf_e1.paragraphs[0].text = "• Выдача корпоративного Wi-Fi (`WLAN-WORKNET`) за 2 секунды."
    tf_e1.paragraphs[0].font.size = Pt(11.5)
    tf_e1.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_e1.add_paragraph()
    p.text = "• Нормализация ФИО/инициалов и поиск учетной записи в домене."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_e1.add_paragraph()
    p.text = "• Single-DC Affinity: исключение рассинхронизации репликации между DC."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s10, 4.8, top_pos, col_w, 5.1, "2. Оргтехника (WinRM / SMB)", border_color=COLOR_ACCENT_PURPLE)
    tb_e2 = s10.shapes.add_textbox(Inches(5.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_e2 = tb_e2.text_frame
    tf_e2.word_wrap = True
    tf_e2.paragraphs[0].text = "• Удаленная установка драйверов принтеров без визита к рабочему месту."
    tf_e2.paragraphs[0].font.size = Pt(11.5)
    tf_e2.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_e2.add_paragraph()
    p.text = "• База моделей оргтехники и автоматическое сопоставление сетевых портов."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_e2.add_paragraph()
    p.text = "• WMI bootstrap и Fail-Fast диагностика."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s10, 8.8, top_pos, col_w, 5.1, "3. Verified Execution Only", border_color=COLOR_ACCENT_BLUE)
    tb_e3 = s10.shapes.add_textbox(Inches(9.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_e3 = tb_e3.text_frame
    tf_e3.word_wrap = True
    tf_e3.paragraphs[0].text = "• Принцип подтвержденного действия:\nСтатус тикета 'Выполнена' выставляется ТОЛЬКО после получения успешного кода возврата от домена."
    tf_e3.paragraphs[0].font.size = Pt(11.5)
    tf_e3.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_e3.add_paragraph()
    p.text = "• Автоматическое списание трудозатрат и протоколирование аудита."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    s10.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Главная сила IntraLink — реальное исполнение. Выдача доступа в Wi-Fi группу Active Directory, "
        "установка принтера по WinRM выполняются в 1 клик. При этом действует принцип Verified Execution: "
        "тикет никогда не закроется, если скрипт в домене вернул ошибку»."
    )

    # =========================================================================
    # СЛАЙД 11: Канбан-пайплайн (Live Showcase 4)
    # =========================================================================
    s11 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s11)
    add_header(s11, "Процессное управление", "Канбан-пайплайн и гибкое ведение очереди", "Визуализация жизненного цикла заявок по стадиям обработки 1-й линии")

    add_card(s11, 0.8, 1.8, 4.8, 5.1, "Преимущества Канбан", border_color=COLOR_ACCENT_EMERALD)
    tb_p11 = s11.shapes.add_textbox(Inches(1.05), Inches(2.4), Inches(4.3), Inches(4.3))
    tf_p11 = tb_p11.text_frame
    tf_p11.word_wrap = True
    p = tf_p11.paragraphs[0]
    p.text = "• Колонки жизненного цикла:\nНовые (Входящие) → В работе (Назначены) → Требует уточнения (Ожидание заявителя) → Выполнены."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p11.add_paragraph()
    p.text = "• Визуальный контроль завалов:\nРуководитель группы 1-й линии сразу видит, на какой стадии скапливаются тикеты."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_p11.add_paragraph()
    p.text = "• Быстрый фокус:\nИнженер кликает на карточку и моментально переходит к инспектору и решению проблемы."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_screenshot_frame(s11, 5.8, 1.8, 6.7, 5.1, "05_kanban_pipeline.png", "Канбан-доска стадий обработки очереди обращений Helpdesk")

    s11.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Помимо классического списка, доступен визуальный Канбан-пайплайн. "
        "Он позволяет дежурному инженеру и тимлиду наглядно оценивать распределение нагрузки "
        "и исключать зависание заявок в ожидании»."
    )

    # =========================================================================
    # СЛАЙД 12: Безопасность и комплаенс
    # =========================================================================
    s12 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s12)
    add_header(s12, "Безопасность и комплаенс", "Enterprise-Grade безопасность и защита данных", "100% On-Premise, криптографическая защита и строгий аудит действий")

    top_pos = 1.8
    add_card(s12, 0.8, top_pos, 3.65, 5.1, "1. Изолированный контур", border_color=COLOR_ACCENT_EMERALD)
    tb_s1 = s12.shapes.add_textbox(Inches(1.05), Inches(top_pos + 0.6), Inches(3.15), Inches(4.1))
    tf_s1 = tb_s1.text_frame
    tf_s1.word_wrap = True
    tf_s1.paragraphs[0].text = "• Полный On-Premise хостинг (Docker, локальный Postgres 16, Redis 7, Core API)."
    tf_s1.paragraphs[0].font.size = Pt(11.5)
    tf_s1.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_s1.add_paragraph()
    p.text = "• Никакие персональные данные, пароли или тексты заявок не отправляются в публичные облачные LLM."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s12, 4.8, top_pos, 3.65, 5.1, "2. Криптография Fernet", border_color=COLOR_ACCENT_BLUE)
    tb_s2 = s12.shapes.add_textbox(Inches(5.05), Inches(top_pos + 0.6), Inches(3.15), Inches(4.1))
    tf_s2 = tb_s2.text_frame
    tf_s2.word_wrap = True
    tf_s2.paragraphs[0].text = "• Доменные учетные данные шифруются симметричным ключом Fernet."
    tf_s2.paragraphs[0].font.size = Pt(11.5)
    tf_s2.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_s2.add_paragraph()
    p.text = "• Бот не хранит пароли пользователей локально."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_s2.add_paragraph()
    p.text = "• Защита от Timing Attacks при валидации ключей (secrets.compare_digest)."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s12, 8.8, top_pos, 3.65, 5.1, "3. Строгий аудит и SLA", border_color=COLOR_ACCENT_PURPLE)
    tb_s3 = s12.shapes.add_textbox(Inches(9.05), Inches(top_pos + 0.6), Inches(3.15), Inches(4.1))
    tf_s3 = tb_s3.text_frame
    tf_s3.word_wrap = True
    tf_s3.paragraphs[0].text = "• Двойное логирование исполнителей: основной инженер + ассистент."
    tf_s3.paragraphs[0].font.size = Pt(11.5)
    tf_s3.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_s3.add_paragraph()
    p.text = "• Фиксация каждого изменения в истории IntraService с точным списанием трудозатрат."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY
    p = tf_s3.add_paragraph()
    p.text = "• Полная прозрачность для службы информационной безопасности."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    s12.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Для руководства и ИБ принципиален вопрос безопасности: система работает строго внутри периметра, "
        "все пароли зашифрованы криптографическим алгоритмом Fernet, а каждое действие логируется с фиксацией времени и исполнителя»."
    )

    # =========================================================================
    # СЛАЙД 13: Бизнес-метрики и ROI
    # =========================================================================
    s13 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s13)
    add_header(s13, "Бизнес-эффект и ROI", "Измеримые результаты внедрения IntraLink AI", "Сравнение ключевых показателей эффективности работы Helpdesk до и после")

    m_w = 2.7
    m_h = 5.1
    
    add_card(s13, 0.8, 1.8, m_w, m_h, "Время первой реакции", border_color=COLOR_ACCENT_BLUE)
    tb_m1 = s13.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(m_w - 0.4), Inches(4.1))
    tf_m1 = tb_m1.text_frame
    tf_m1.word_wrap = True
    p = tf_m1.paragraphs[0]
    p.text = "-96%"
    p.font.name = FONT_MAIN
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_BLUE
    p = tf_m1.add_paragraph()
    p.text = "\nБыло: 15–20 мин\nСтало: 30 секунд\n\nМгновенные Push-уведомления в Telegram через Redis Streams."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s13, 3.8, 1.8, m_w, m_h, "MTTR (Время решения)", border_color=COLOR_ACCENT_EMERALD)
    tb_m2 = s13.shapes.add_textbox(Inches(4.0), Inches(2.5), Inches(m_w - 0.4), Inches(4.1))
    tf_m2 = tb_m2.text_frame
    tf_m2.word_wrap = True
    p = tf_m2.paragraphs[0]
    p.text = "-94%"
    p.font.name = FONT_MAIN
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_EMERALD
    p = tf_m2.add_paragraph()
    p.text = "\nБыло: 25 минут\nСтало: 1.5 минуты\n\nДля типовых заявок (Wi-Fi, пароли, принтеры, консультации)."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s13, 6.8, 1.8, m_w, m_h, "Автоматизация 1-й линии", border_color=COLOR_ACCENT_PURPLE)
    tb_m3 = s13.shapes.add_textbox(Inches(7.0), Inches(2.5), Inches(m_w - 0.4), Inches(4.1))
    tf_m3 = tb_m3.text_frame
    tf_m3.word_wrap = True
    p = tf_m3.paragraphs[0]
    p.text = "до 45%"
    p.font.name = FONT_MAIN
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_PURPLE
    p = tf_m3.add_paragraph()
    p.text = "\nОчереди обрабатывается в полуавтоматическом и автоматическом режиме без ручного набора текста."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s13, 9.8, 1.8, m_w, m_h, "Экономия времени", border_color=COLOR_ACCENT_AMBER)
    tb_m4 = s13.shapes.add_textbox(Inches(10.0), Inches(2.5), Inches(m_w - 0.4), Inches(4.1))
    tf_m4 = tb_m4.text_frame
    tf_m4.word_wrap = True
    p = tf_m4.paragraphs[0]
    p.text = "160 ч/мес"
    p.font.name = FONT_MAIN
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT_AMBER
    p = tf_m4.add_paragraph()
    p.text = "\nВысвобождается до 1 FTE инженера для решения сложных инфраструктурных задач развития."
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_TEXT_SECONDARY

    s13.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Главный итог в цифрах: время первичной реакции упало на 96% (с 15 минут до 30 секунд), "
        "время закрытия типовых тикетов сократилось с 25 минут до полутора минут, а суммарная экономия "
        "достигает 160 человеко-часов в месяц, высвобождая инженера для системных проектов»."
    )

    # =========================================================================
    # СЛАЙД 14: Roadmap развития
    # =========================================================================
    s14 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s14)
    add_header(s14, "Стратегия развития", "Дорожная карта развития IntraLink AI", "Ключевые направления масштабирования платформы на следующие кварталы")

    col_w = 3.65
    top_pos = 1.8
    add_card(s14, 0.8, top_pos, col_w, 5.1, "Этап 1: Реализовано (Внедрено)", border_color=COLOR_ACCENT_EMERALD)
    tb_rd1 = s14.shapes.add_textbox(Inches(1.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_rd1 = tb_rd1.text_frame
    tf_rd1.word_wrap = True
    tf_rd1.paragraphs[0].text = "✓ FastAPI Core Gateway & Web UI\n✓ RAG база знаний pgvector + FastEmbed\n✓ Пакетный триаж очереди (/triage)\n✓ Авто-дубликаты и авто-редирект\n✓ Роботизация AD WLAN & WinRM\n✓ Локальная LLM суммаризация Qwen2.5"
    tf_rd1.paragraphs[0].font.size = Pt(11.5)
    tf_rd1.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s14, 4.8, top_pos, col_w, 5.1, "Этап 2: Ближайший квартал (Q3-Q4)", border_color=COLOR_ACCENT_BLUE)
    tb_rd2 = s14.shapes.add_textbox(Inches(5.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_rd2 = tb_rd2.text_frame
    tf_rd2.word_wrap = True
    tf_rd2.paragraphs[0].text = "• Интерактивный AI-саппорт в Telegram для пользователей (самообслуживание).\n• Аудио-транскрипция голосовых заявок и звонков (Whisper local).\n• Расширение библиотеки драйверов оргтехники и SNMP-мониторинг остатка картриджей."
    tf_rd2.paragraphs[0].font.size = Pt(11.5)
    tf_rd2.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s14, 8.8, top_pos, col_w, 5.1, "Этап 3: Стратегическая перспектива", border_color=COLOR_ACCENT_PURPLE)
    tb_rd3 = s14.shapes.add_textbox(Inches(9.05), Inches(top_pos + 0.6), Inches(col_w - 0.5), Inches(4.1))
    tf_rd3 = tb_rd3.text_frame
    tf_rd3.word_wrap = True
    tf_rd3.paragraphs[0].text = "• Предиктивная аналитика сбоев рабочих станций (предупреждение сбоев до создания тикета).\n• Автоматический аудит узких мест ИТ-процессов.\n• Кросс-сервисная интеграция с 1C и корпоративной телефонией."
    tf_rd3.paragraphs[0].font.size = Pt(11.5)
    tf_rd3.paragraphs[0].font.color.rgb = COLOR_TEXT_SECONDARY

    s14.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«Наш роадмап включает 3 этапа. Фундамент уже работает: RAG, триаж, AD и локальная LLM внедрены. "
        "В ближайших планах — бот самообслуживания для пользователей и распознавание голосовых сообщений. "
        "В стратегической перспективе — предиктивное обслуживание до обращения заявителя»."
    )

    # =========================================================================
    # СЛАЙД 15: Заключение и Q&A
    # =========================================================================
    s15 = prs.slides.add_slide(blank_slide_layout)
    add_bg(s15)
    
    tb_fin = s15.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(2.5))
    tf_fin = tb_fin.text_frame
    tf_fin.word_wrap = True
    p = tf_fin.paragraphs[0]
    p.text = "IntraLink AI"
    p.font.name = FONT_MAIN
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_PRIMARY

    p2 = tf_fin.add_paragraph()
    p2.text = "Интеллект. Скорость. Безопасность."
    p2.font.name = FONT_MAIN
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_ACCENT_EMERALD

    p3 = tf_fin.add_paragraph()
    p3.text = "\nГотовы ответить на ваши вопросы и продемонстрировать систему в действии."
    p3.font.name = FONT_MAIN
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEXT_SECONDARY

    add_card(s15, 0.8, 4.6, 11.7, 2.0, "Контакты и ресурсы системы", border_color=COLOR_ACCENT_BLUE)
    tb_fc = s15.shapes.add_textbox(Inches(1.05), Inches(5.15), Inches(11.2), Inches(1.2))
    tf_fc = tb_fc.text_frame
    tf_fc.word_wrap = True
    p = tf_fc.paragraphs[0]
    p.text = "• Веб-панель управления Helpdesk: http://localhost:8000/admin\n• Интерактивная Web-версия презентации: /docs/presentations/index.html\n• Архитектура и спецификации: /docs/architecture.md"
    p.font.name = FONT_MONO
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_TEXT_PRIMARY

    s15.notes_slide.notes_text_frame.text = (
        "Скрипт для спикера:\n"
        "«IntraLink AI превратил 1-ю линию поддержки из перегруженного узкого горлышка "
        "в прозрачный, защищенный и быстрый сервис. Спасибо за внимание! "
        "Готов ответить на вопросы и провести живую демонстрацию системы»."
    )

    output_dir = os.path.join("docs", "presentations")
    os.makedirs(output_dir, exist_ok=True)
    
    # Сохраняем в основной файл и в latest
    primary_path = os.path.join(output_dir, "IntraLink_AI_Executive_Presentation.pptx")
    latest_path = os.path.join(output_dir, "IntraLink_AI_Presentation_Latest.pptx")
    
    saved_paths = []
    try:
        prs.save(primary_path)
        saved_paths.append(primary_path)
        print(f"Презентация сохранена в: {primary_path}")
    except PermissionError:
        print(f"Файл {primary_path} занят другим процессом (например, открыт в PowerPoint).")
    
    try:
        prs.save(latest_path)
        saved_paths.append(latest_path)
        print(f"Презентация сохранена в: {latest_path}")
    except Exception as e:
        print(f"Ошибка сохранения {latest_path}: {e}")

    return saved_paths

if __name__ == "__main__":
    create_presentation()
